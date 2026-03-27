import os
import re
import base64
from typing import Dict, Any, List
from app.skills.base import BaseProcessor, ProcessorRegistry
from app.core.constants import DocumentConstants
import PyPDF2
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image
import logging

logger = logging.getLogger(__name__)

# 尝试导入OCR库
try:
    import pytesseract
    PYTESSERACT_AVAILABLE = True
except ImportError:
    PYTESSERACT_AVAILABLE = False
    logger.warning("pytesseract 未安装，图片OCR功能不可用")


@ProcessorRegistry.register("DocumentParser")
class DocumentParser(BaseProcessor):
    
    @property
    def name(self) -> str:
        return "DocumentParser"
    
    async def process(self, context: Dict[str, Any], params: Dict[str, Any]) -> Dict[str, Any]:
        file_path = context.get("file_path")
        # Check both original_filename and document_name for compatibility
        original_filename = context.get("original_filename") or context.get("document_name")
        
        if not file_path:
            raise ValueError("file_path is required")
        
        # Use original filename if provided, otherwise use the file_path basename
        display_filename = original_filename if original_filename else os.path.basename(file_path)
        
        file_ext = os.path.splitext(file_path)[1].lower().replace(".", "")
        supported_formats = params.get("supported_formats", DocumentConstants.SUPPORTED_FORMATS)
        
        if file_ext not in supported_formats:
            raise ValueError(f"Unsupported format: {file_ext}")
        
        text = ""
        if file_ext == "pdf":
            text = self._parse_pdf(file_path, display_filename)
        elif file_ext == "docx":
            text = self._parse_docx(file_path)
        elif file_ext == "txt":
            text = self._parse_txt(file_path)
        elif file_ext == "xlsx":
            text = self._parse_xlsx(file_path, display_filename)
        elif file_ext == "pptx":
            text = self._parse_pptx(file_path, display_filename)
        elif file_ext in ["png", "jpg", "jpeg", "gif", "bmp", "webp"]:
            text = self._parse_image(file_path, display_filename)
        
        return {
            "document_text": text,
            "document_format": file_ext,
            "document_id": str(os.path.basename(file_path)),
            "document_name": display_filename
        }
    
    def _parse_pdf(self, file_path: str, display_filename: str = None) -> str:
        text_parts = []
        # Use display_filename if provided, otherwise fall back to file_path basename
        filename = display_filename if display_filename else os.path.basename(file_path)
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(f"[{filename} - Page {i+1}]\n{page_text}")
        return "\n\n".join(text_parts)
    
    def _parse_docx(self, file_path: str) -> str:
        doc = DocxDocument(file_path)
        return "\n\n".join([para.text for para in doc.paragraphs if para.text])
    
    def _parse_txt(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    
    def _parse_xlsx(self, file_path: str, display_filename: str = None) -> str:
        """
        解析 Excel 文件，确保每行数据的完整性
        
        关键设计：
        1. 每行数据作为一个完整的文本块，包含表头信息
        2. 使用 "键: 值" 格式，增强语义关联
        3. 确保人员、上级等关联信息在同一文本块中
        """
        filename = display_filename if display_filename else os.path.basename(file_path)
        text_parts = []
        
        wb = load_workbook(file_path, data_only=True)
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"[{filename} - 工作表: {sheet_name}]")
            
            rows = list(sheet.iter_rows(values_only=True))
            if not rows:
                continue
            
            # 第一行作为表头
            headers = [str(cell) if cell is not None else f"列{i+1}" for i, cell in enumerate(rows[0])]
            text_parts.append(f"表头: {' | '.join(headers)}")
            text_parts.append("")
            
            # 处理数据行 - 每行作为一个完整的记录块
            for row_idx, row in enumerate(rows[1:], 2):
                if not any(cell is not None for cell in row):
                    continue
                
                # 构建键值对格式的行记录
                row_data = []
                for i, (header, cell) in enumerate(zip(headers, row)):
                    if cell is not None:
                        row_data.append(f"{header}: {cell}")
                
                if row_data:
                    # 将整行数据作为一个完整的文本块
                    # 格式: 记录 N: 姓名: xxx | 上级: xxx | 部门: xxx
                    record_text = f"记录 {row_idx-1}: " + " | ".join(row_data)
                    text_parts.append(record_text)
            
            text_parts.append("")
        
        return "\n".join(text_parts)
    
    def _parse_pptx(self, file_path: str, display_filename: str = None) -> str:
        filename = display_filename if display_filename else os.path.basename(file_path)
        text_parts = []
        
        prs = Presentation(file_path)
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            text_parts.append(f"[{filename} - 幻灯片 {slide_idx}]")
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
                
                if shape.has_table:
                    table = shape.table
                    table_text = []
                    for row in table.rows:
                        row_text = " | ".join([cell.text for cell in row.cells])
                        table_text.append(row_text)
                    text_parts.append("表格:\n" + "\n".join(table_text))
            
            text_parts.append("")
        
        return "\n".join(text_parts)
    
    def _parse_image(self, file_path: str, display_filename: str = None) -> str:
        filename = display_filename if display_filename else os.path.basename(file_path)
        
        try:
            img = Image.open(file_path)
            width, height = img.size
            format_name = img.format if img.format else "Unknown"
            mode = img.mode
            
            text_parts = [
                f"[图片文件: {filename}]",
                f"格式: {format_name}",
                f"尺寸: {width} x {height} 像素",
                f"颜色模式: {mode}",
            ]
            
            if hasattr(img, 'info'):
                exif_info = []
                if 'exif' in img.info:
                    exif_info.append("包含EXIF元数据")
                if 'dpi' in img.info:
                    exif_info.append(f"DPI: {img.info['dpi']}")
                if exif_info:
                    text_parts.append("元数据: " + ", ".join(exif_info))
            
            # 尝试OCR识别文字
            ocr_text = self._extract_text_from_image(img)
            if ocr_text:
                text_parts.append("")
                text_parts.append("[图片中的文字内容]")
                text_parts.append(ocr_text)
            else:
                text_parts.append("")
                text_parts.append("注意: 未从图片中识别到文字内容。")
            
            return "\n".join(text_parts)
            
        except Exception as e:
            logger.error(f"解析图片失败: {file_path}, 错误: {e}")
            return f"[图片文件: {filename}]\n解析失败: {str(e)}"
    
    def _extract_text_from_image(self, img: Image.Image) -> str:
        """使用OCR从图片中提取文字"""
        if not PYTESSERACT_AVAILABLE:
            logger.warning("pytesseract 未安装，跳过OCR")
            return ""
        
        try:
            # 转换为RGB模式（如果是RGBA或其他模式）
            if img.mode != 'RGB':
                img = img.convert('RGB')
            
            # 使用pytesseract进行OCR识别
            # lang='chi_sim+eng' 表示使用中文简体和英文
            text = pytesseract.image_to_string(img, lang='chi_sim+eng')
            
            # 清理文本
            text = text.strip()
            
            if text:
                logger.info(f"OCR识别成功，提取了 {len(text)} 个字符")
                return text
            else:
                logger.info("OCR未识别到文字")
                return ""
                
        except Exception as e:
            logger.error(f"OCR识别失败: {e}")
            return ""


@ProcessorRegistry.register("TextSplitter")
class TextSplitter(BaseProcessor):
    
    @property
    def name(self) -> str:
        return "TextSplitter"
    
    async def process(self, context: Dict[str, Any], params: Dict[str, Any]) -> Dict[str, Any]:
        text = context.get("document_text", "")
        document_format = context.get("document_format", "")
        chunk_size = params.get("chunk_size", 500)
        chunk_overlap = params.get("chunk_overlap", 50)
        
        # 对于结构化数据（xlsx），使用特殊处理策略
        if document_format in ["xlsx", "pptx"]:
            chunks = self._split_structured_data(text, chunk_size)
        else:
            chunks = self._split_text(text, chunk_size, chunk_overlap)
        
        return {
            "chunks": chunks,
            "chunk_count": len(chunks),
            "document_id": context.get("document_id"),
            "document_name": context.get("document_name")
        }
    
    def _split_text(self, text: str, chunk_size: int, overlap: int) -> list:
        paragraphs = text.split("\n\n")
        chunks = []
        current_chunk = ""
        
        for para in paragraphs:
            if len(current_chunk) + len(para) <= chunk_size:
                current_chunk += para + "\n\n"
            else:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                current_chunk = para + "\n\n"
        
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        return chunks
    
    def _split_structured_data(self, text: str, chunk_size: int) -> list:
        """
        处理结构化数据（如 Excel、PPT）
        
        关键原则：
        1. 保持每行记录的完整性，不将一行拆分到多个 chunk
        2. 每个 chunk 包含表头信息，便于理解上下文
        3. 按行分割，确保检索时能获取完整记录
        """
        lines = text.split("\n")
        chunks = []
        
        # 提取表头信息（通常在文件开头）
        header_lines = []
        data_lines = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            # 表头行特征：包含 "表头:" 或在工作表标题之后
            if "表头:" in line or "工作表:" in line or line.startswith("["):
                header_lines.append(line)
            else:
                data_lines.append(line)
        
        # 构建 chunk：每个 chunk 包含表头 + 若干完整记录
        header_text = "\n".join(header_lines)
        current_chunk = header_text + "\n\n" if header_text else ""
        
        for line in data_lines:
            # 如果当前行是记录行（以 "记录" 开头），单独处理
            if line.startswith("记录 "):
                # 如果当前 chunk 加上新行会超过限制，保存当前 chunk
                if len(current_chunk) + len(line) > chunk_size and current_chunk.strip() != header_text.strip():
                    chunks.append(current_chunk.strip())
                    current_chunk = header_text + "\n\n" + line + "\n"
                else:
                    current_chunk += line + "\n"
            else:
                # 其他行（如表头说明）
                current_chunk += line + "\n"
        
        if current_chunk.strip() and current_chunk.strip() != header_text.strip():
            chunks.append(current_chunk.strip())
        
        # 如果没有生成任何 chunk（可能是空文件），返回原始文本
        if not chunks and text.strip():
            chunks = [text.strip()]
        
        return chunks


@ProcessorRegistry.register("SmartTextSplitter")
class SmartTextSplitter(BaseProcessor):
    
    CHUNK_CONFIGS = {
        "policy": {
            "chunk_size": 400,
            "chunk_overlap": 80,
            "separators": ["\n\n", "\n", "。", "；", "，", " "],
            "description": "政策文档 - 保持条款完整性"
        },
        "technical": {
            "chunk_size": 600,
            "chunk_overlap": 100,
            "separators": ["\n\n", "\n", "```", "。", "；", " "],
            "description": "技术文档 - 保留代码块完整"
        },
        "faq": {
            "chunk_size": 300,
            "chunk_overlap": 50,
            "separators": ["\n\n", "\nQ:", "\nA:", "。"],
            "description": "FAQ文档 - 问答对不拆分"
        },
        "general": {
            "chunk_size": 500,
            "chunk_overlap": 50,
            "separators": ["\n\n", "\n", "。", "！", "？", "；", "，", " "],
            "description": "通用文档 - 平衡策略"
        }
    }
    
    @property
    def name(self) -> str:
        return "SmartTextSplitter"
    
    async def process(self, context: Dict[str, Any], params: Dict[str, Any]) -> Dict[str, Any]:
        text = context.get("document_text", "")
        doc_type = params.get("doc_type", "general")
        custom_config = params.get("chunk_config")
        
        config = custom_config or self.CHUNK_CONFIGS.get(doc_type, self.CHUNK_CONFIGS["general"])
        
        chunks = self._split_text_smart(
            text,
            chunk_size=config.get("chunk_size", 500),
            chunk_overlap=config.get("chunk_overlap", 50),
            separators=config.get("separators", ["\n\n", "\n", "。", " "])
        )
        
        chunks_with_meta = self._add_chunk_metadata(chunks, doc_type)
        
        logger.info(f"智能分块完成: 文档类型={doc_type}, 分块数={len(chunks)}")
        
        return {
            "chunks": chunks_with_meta,
            "chunk_count": len(chunks),
            "doc_type": doc_type,
            "config": config
        }
    
    def _split_text_smart(
        self, 
        text: str, 
        chunk_size: int, 
        chunk_overlap: int,
        separators: List[str]
    ) -> List[str]:
        chunks = []
        
        text = self._preprocess_text(text)
        
        sections = self._split_by_separators(text, separators)
        
        current_chunk = ""
        for section in sections:
            if len(section) > chunk_size:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                    current_chunk = ""
                
                sub_chunks = self._split_large_section(section, chunk_size, chunk_overlap)
                chunks.extend(sub_chunks)
            elif len(current_chunk) + len(section) <= chunk_size:
                current_chunk += section
            else:
                if current_chunk:
                    chunks.append(current_chunk.strip())
                current_chunk = section
        
        if current_chunk:
            chunks.append(current_chunk.strip())
        
        if chunk_overlap > 0 and len(chunks) > 1:
            chunks = self._add_overlap(chunks, chunk_overlap)
        
        return [c for c in chunks if c.strip()]
    
    def _preprocess_text(self, text: str) -> str:
        text = re.sub(r'\n{3,}', '\n\n', text)
        text = re.sub(r' {2,}', ' ', text)
        return text.strip()
    
    def _split_by_separators(self, text: str, separators: List[str]) -> List[str]:
        if not separators:
            return [text]
        
        primary_sep = separators[0]
        sections = text.split(primary_sep)
        
        if len(sections) == 1 and len(separators) > 1:
            return self._split_by_separators(text, separators[1:])
        
        return [s.strip() for s in sections if s.strip()]
    
    def _split_large_section(self, section: str, chunk_size: int, overlap: int) -> List[str]:
        chunks = []
        start = 0
        
        while start < len(section):
            end = start + chunk_size
            
            if end < len(section):
                break_point = self._find_break_point(section, end)
                end = break_point
            
            chunk = section[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            start = end - overlap if end < len(section) else end
        
        return chunks
    
    def _find_break_point(self, text: str, position: int) -> int:
        break_chars = ["。", "！", "？", "；", "\n", "，", " "]
        
        for i in range(position, max(position - 100, 0), -1):
            if text[i] in break_chars:
                return i + 1
        
        return position
    
    def _add_overlap(self, chunks: List[str], overlap: int) -> List[str]:
        if len(chunks) <= 1:
            return chunks
        
        overlapped = [chunks[0]]
        
        for i in range(1, len(chunks)):
            prev_chunk = chunks[i-1]
            overlap_text = prev_chunk[-overlap:] if len(prev_chunk) > overlap else prev_chunk
            overlapped.append(overlap_text + chunks[i])
        
        return overlapped
    
    def _add_chunk_metadata(self, chunks: List[str], doc_type: str) -> List[Dict[str, Any]]:
        return [
            {
                "content": chunk,
                "metadata": {
                    "chunk_index": i,
                    "doc_type": doc_type,
                    "char_count": len(chunk)
                }
            }
            for i, chunk in enumerate(chunks)
        ]
    
    def detect_doc_type(self, text: str, filename: str = "") -> str:
        if "第" in text and "条" in text and ("规定" in text or "应当" in text):
            return "policy"
        
        if "```" in text or "def " in text or "function " in text:
            return "technical"
        
        if re.search(r'Q[:：]|问[:：]', text) and re.search(r'A[:：]|答[:：]', text):
            return "faq"
        
        if filename:
            ext = os.path.splitext(filename)[1].lower()
            if ext in [".pdf"]:
                return "policy"
            elif ext in [".md", ".rst"]:
                return "technical"
        
        return "general"
